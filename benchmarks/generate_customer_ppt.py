import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except Exception as exc:
    print(f"ERROR: python-pptx import failed: {exc}", file=sys.stderr)
    sys.exit(1)

ROOT = Path(__file__).resolve().parent
CUSTOMER_DIR = ROOT / "customer_deck"
OUT_DIR = ROOT / "customer_ppt"

TITLE_COLOR = "0F172A"
ACCENT_COLOR = "0F766E"
TEXT_COLOR = "111827"
MUTED_COLOR = "475569"
LIGHT_BG = "F0FDFA"


def load_json(path: Path, default):
    if not path.exists():
        return default
    return json.loads(path.read_text(encoding="utf-8-sig"))


def latest_file(dir_path: Path, pattern: str):
    files = sorted(dir_path.glob(pattern))
    if not files:
        raise RuntimeError(f"No files found for pattern: {pattern}")
    return files[-1]


def set_rgb(font_obj, hex_rgb: str):
    font_obj.color.rgb = RGBColor.from_string(hex_rgb)


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.6), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(28)
    set_rgb(r.font, TITLE_COLOR)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(11.0), Inches(0.45))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(13)
        set_rgb(r2.font, MUTED_COLOR)


def add_bullets(slide, bullets, left=0.9, top=1.6, width=10.8, height=5.2, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(10)
        p.font.size = Pt(font_size)
        set_rgb(p.font, TEXT_COLOR)


def add_metric_row(slide, metrics):
    x_positions = [0.7, 3.2, 5.7, 8.2, 10.4]
    widths = [2.2, 2.2, 2.2, 1.8, 1.8]
    titles = [
        ("Retrieval Hit", metrics.get("retrieval_hit_rate")),
        ("Soft Hit", metrics.get("soft_hit_rate")),
        ("Context Reduction", metrics.get("context_reduction_percent")),
        ("Pruned Lines", metrics.get("total_pruned_lines")),
        ("Unbeaten", f"{metrics.get('unbeaten_scenarios')} / {metrics.get('scenario_count')}")
    ]
    for i, (label, value) in enumerate(titles):
        shape = slide.shapes.add_shape(1, Inches(x_positions[i]), Inches(1.75), Inches(widths[i]), Inches(1.25))
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(LIGHT_BG)
        shape.line.color.rgb = RGBColor.from_string("99F6E4")
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = label
        r1.font.size = Pt(11)
        r1.font.bold = True
        set_rgb(r1.font, MUTED_COLOR)

        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = str(value)
        r2.font.size = Pt(18)
        r2.font.bold = True
        set_rgb(r2.font, ACCENT_COLOR)


def add_comparison_table(slide, table_like):
    rows = len(table_like) + 1
    cols = 5
    table = slide.shapes.add_table(rows, cols, Inches(0.55), Inches(1.55), Inches(11.8), Inches(4.8)).table
    headers = ["Mode", "Retrieval", "Soft", "Context Reduction", "Removed"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string("CCFBF1")
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True

    for r, (mode, vals) in enumerate(table_like.items(), start=1):
        table.cell(r, 0).text = mode
        table.cell(r, 1).text = str(vals.get("retrieval_hit_rate"))
        table.cell(r, 2).text = str(vals.get("soft_hit_rate"))
        table.cell(r, 3).text = str(vals.get("context_reduction_percent"))
        table.cell(r, 4).text = str(vals.get("repeated_explanation_items_removed"))
        for c in range(cols):
            p = table.cell(r, c).text_frame.paragraphs[0]
            p.font.size = Pt(11)


def add_footer(slide, run_id):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(6.95), Inches(6.5), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"Run ID: {run_id}"
    r.font.size = Pt(9)
    set_rgb(r.font, MUTED_COLOR)


def build_slide(prs, slide_data, run_id):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, slide_data.get("title", ""), slide_data.get("subtitle"))

    if slide_data.get("slide_number") == 1:
        add_bullets(slide, slide_data.get("bullets", []), top=1.9, height=3.2, font_size=22)
    elif slide_data.get("title") == "Core Metrics":
        metric_map = {}
        for bullet in slide_data.get("bullets", []):
            if ":" in bullet:
                k, v = bullet.split(":", 1)
                metric_map[k.strip().lower().replace(" ", "_")] = v.strip()
        unbeaten = metric_map.get("unbeaten_scenarios", "")
        unbeaten_left = unbeaten.split("/")[0].strip() if "/" in unbeaten else unbeaten
        unbeaten_right = unbeaten.split("/")[1].strip() if "/" in unbeaten else None
        add_metric_row(slide, {
            "retrieval_hit_rate": metric_map.get("retrieval_hit_rate"),
            "soft_hit_rate": metric_map.get("soft_hit_rate"),
            "context_reduction_percent": metric_map.get("context_reduction_percent"),
            "total_pruned_lines": metric_map.get("total_pruned_lines"),
            "unbeaten_scenarios": unbeaten_left,
            "scenario_count": unbeaten_right,
        })
        add_bullets(slide, [
            "High continuity retention.",
            "Lower context burden.",
            "Strong performance under robustness validation."
        ], top=3.4, height=2.5, font_size=18)
    elif "table_like" in slide_data:
        add_comparison_table(slide, slide_data["table_like"])
    else:
        add_bullets(slide, slide_data.get("bullets", []), top=1.7, height=5.3, font_size=18)

    add_footer(slide, run_id)


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    summary_path = latest_file(CUSTOMER_DIR, "customer-deck-summary-*.json")
    summary = load_json(summary_path, {})
    if not summary:
        raise RuntimeError("Could not load customer deck summary.")

    run_id = summary.get("run_id", "unknown-run")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in summary.get("slides", []):
        build_slide(prs, slide_data, run_id)

    pptx_path = OUT_DIR / f"portable-memory-customer-deck-{run_id}.pptx"
    prs.save(str(pptx_path))

    print("Customer PowerPoint generated.")
    print(f"Source summary: {summary_path}")
    print(f"PPTX: {pptx_path}")


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        sys.exit(1)
